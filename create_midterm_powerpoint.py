#!/usr/bin/env python3
"""
Create PowerPoint presentation with midterm questions, answers, and code examples.
Perfect for open-book database exam reference.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame

    for item in content_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(14)

    return slide

def add_code_slide(prs, title, code_text):
    """Add a slide with code."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)

    # Add code box
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(6))
    code_frame = code_box.text_frame
    code_frame.word_wrap = True

    code_para = code_frame.paragraphs[0]
    code_para.text = code_text
    code_para.font.name = 'Courier New'
    code_para.font.size = Pt(9)

    return slide

def add_qa_slide(prs, question, answer, code=""):
    """Add a question and answer slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Question
    q_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.2))
    q_frame = q_box.text_frame
    q_para = q_frame.paragraphs[0]
    q_para.text = f"Q: {question}"
    q_para.font.size = Pt(16)
    q_para.font.bold = True
    q_para.font.color.rgb = RGBColor(0, 102, 204)
    q_frame.word_wrap = True

    # Answer
    a_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(0.8))
    a_frame = a_box.text_frame
    a_para = a_frame.paragraphs[0]
    a_para.text = f"A: {answer}"
    a_para.font.size = Pt(12)
    a_frame.word_wrap = True

    # Code (if provided)
    if code:
        code_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(4.5))
        code_frame = code_box.text_frame
        code_frame.word_wrap = True
        code_para = code_frame.paragraphs[0]
        code_para.text = code
        code_para.font.name = 'Courier New'
        code_para.font.size = Pt(8)

    return slide

def create_comprehensive_powerpoint():
    """Create comprehensive PowerPoint with all midterm content."""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== COVER SLIDE =====
    add_title_slide(prs,
                   "Advanced Database\nMidterm Study Guide",
                   "Complete Questions & Solutions\nOpen-Book Exam Reference")

    # ===== TABLE OF CONTENTS =====
    add_content_slide(prs, "Table of Contents", [
        "Section 1: Managing Database Connections",
        "Section 2: Users, Roles & Permissions",
        "Section 3: Tables, Data Types & Arrays",
        "Section 4: Backup & Restore",
        "Quick Reference: Common Commands",
        "Quick Reference: Array Operations",
        "Quick Reference: String Functions"
    ])

    # ========================================================================
    # SECTION 1: DATABASE CONNECTIONS
    # ========================================================================
    add_title_slide(prs, "Section 1", "Managing Database Connections")

    # Q1a: List connections
    add_qa_slide(prs,
        "How can you retrieve a list of recent connections and process IDs (PIDs)?",
        "Use the pg_stat_activity system view to see all active connections, queries, and session information.",
        """-- Basic query to see all connections
SELECT pid, usename, datname, state, query_start, query
FROM pg_stat_activity;

-- See only active connections
SELECT pid, usename, datname, application_name,
       client_addr, state, query
FROM pg_stat_activity
WHERE state = 'active';

-- See connections for specific database
SELECT pid, usename, state, query
FROM pg_stat_activity
WHERE datname = 'your_database';"""
    )

    # Q1b: Cancel and terminate
    add_qa_slide(prs,
        "How can you cancel active queries and terminate connections with a given PID?",
        "Use pg_cancel_backend(pid) to cancel queries gracefully, or pg_terminate_backend(pid) to forcefully close connections.",
        """-- Find the PID
SELECT pid, usename, datname, state, query
FROM pg_stat_activity
WHERE usename = 'problematic_user';

-- Cancel running query (gentle)
SELECT pg_cancel_backend(12345);

-- Terminate connection (forceful)
SELECT pg_terminate_backend(12345);

-- Verify connection is gone
SELECT pid FROM pg_stat_activity WHERE pid = 12345;"""
    )

    # Q1c: Kill all connections
    add_qa_slide(prs,
        "How can you kill all connections belonging to a specific role?",
        "Use a query with pg_terminate_backend() on all PIDs for that role. Always exclude your own PID!",
        """-- Terminate all connections for 'dev_user'
SELECT pg_terminate_backend(pid)
FROM pg_stat_activity
WHERE usename = 'dev_user'
  AND pid <> pg_backend_pid();

-- Terminate all connections to a database
SELECT pg_terminate_backend(pid)
FROM pg_stat_activity
WHERE datname = 'sample_db'
  AND pid <> pg_backend_pid();

-- Verify
SELECT pid, usename FROM pg_stat_activity
WHERE usename = 'dev_user';"""
    )

    # ========================================================================
    # SECTION 2: USERS, ROLES & PERMISSIONS
    # ========================================================================
    add_title_slide(prs, "Section 2", "Users, Roles & Permissions")

    # Q2a: Create user with expiration
    add_qa_slide(prs,
        "Create user 'dev_user' with password, valid until specific date, with read/write on 'employee_data'.",
        "Use CREATE USER with VALID UNTIL clause, then GRANT SELECT, INSERT, UPDATE, DELETE.",
        """-- Create user with expiration
CREATE USER dev_user
WITH PASSWORD 'SecurePass123!'
VALID UNTIL '2025-12-31';

-- Grant read permission
GRANT SELECT ON employee_data TO dev_user;

-- Grant write permissions
GRANT INSERT, UPDATE, DELETE ON employee_data TO dev_user;

-- Or grant all at once:
GRANT SELECT, INSERT, UPDATE, DELETE
ON employee_data TO dev_user;

-- Verify
SELECT usename, valuntil FROM pg_user
WHERE usename = 'dev_user';"""
    )

    # Q2b: Create role with inheritance
    add_qa_slide(prs,
        "Create role 'dev_team', assign 'dev_user' to it, grant access to all tables with inheritance.",
        "Create role with INHERIT, grant privileges on ALL TABLES, use ALTER DEFAULT PRIVILEGES for future tables.",
        """-- Create role with inheritance
CREATE ROLE dev_team WITH INHERIT;

-- Grant access to all current tables
GRANT SELECT, INSERT, UPDATE, DELETE
ON ALL TABLES IN SCHEMA public
TO dev_team;

-- Grant privileges on future tables
ALTER DEFAULT PRIVILEGES IN SCHEMA public
GRANT SELECT, INSERT, UPDATE, DELETE
ON TABLES TO dev_team;

-- Assign user to role
GRANT dev_team TO dev_user;"""
    )

    # Q2c: Add new developer
    add_qa_slide(prs,
        "Add new developer to 'dev_team' role and test access to 'employee_data'.",
        "Create new user with LOGIN, grant role membership, test with SET ROLE.",
        """-- Create new developer
CREATE USER new_developer
WITH PASSWORD 'DevPass456!'
LOGIN;

-- Add to dev_team role
GRANT dev_team TO new_developer;

-- Test access
SET ROLE new_developer;
SELECT * FROM employee_data LIMIT 5;

-- Test write access
INSERT INTO employee_data (name, email)
VALUES ('Test', 'test@example.com');

-- Reset role
RESET ROLE;"""
    )

    # ========================================================================
    # SECTION 3: TABLES, DATA TYPES & ARRAYS
    # ========================================================================
    add_title_slide(prs, "Section 3", "Tables, Data Types & Arrays")

    # Q3a: Create database and table
    add_qa_slide(prs,
        "Create database 'sample_db' with table 'sample_table' (id serial, name varchar(50), age numeric(4,2), description text).",
        "Use CREATE DATABASE, then CREATE TABLE with appropriate data types.",
        """-- Create database
CREATE DATABASE sample_db;

-- Connect to it
\\c sample_db

-- Create table
CREATE TABLE sample_table (
    id SERIAL PRIMARY KEY,
    name VARCHAR(50),
    age NUMERIC(4, 2),
    description TEXT
);

-- Verify structure
\\d sample_table"""
    )

    # Q3b: Load from CSV
    add_qa_slide(prs,
        "Populate 'sample_table' with data from CSV file.",
        "Use \\copy command in psql (client-side) or COPY (server-side).",
        """-- Using \\copy (client-side, no superuser needed)
\\copy sample_table(name, age, description)
FROM 'sample_data.csv' WITH CSV HEADER

-- Using COPY (server-side, needs privileges)
COPY sample_table(name, age, description)
FROM '/path/to/sample_data.csv'
WITH (FORMAT csv, HEADER true);

-- Manual INSERT
INSERT INTO sample_table (name, age, description)
VALUES
('John Smith', 25, 'Description 1'),
('Jane Doe', 33, 'Description 2');"""
    )

    # Q3c/d: String functions
    add_qa_slide(prs,
        "Pad 'name' to 20 chars (left) and trim whitespace from 'description'.",
        "Use LPAD for padding and TRIM for removing whitespace.",
        """-- Left-pad name to 20 characters
SELECT id,
       LPAD(name, 20, ' ') AS padded_name,
       age, description
FROM sample_table;

-- Trim whitespace from description
SELECT id, name, age,
       TRIM(description) AS trimmed,
       LTRIM(description) AS left_trim,
       RTRIM(description) AS right_trim
FROM sample_table;

-- Combined
SELECT LPAD(name, 20, ' ') AS name,
       TRIM(description) AS desc
FROM sample_table;"""
    )

    # Q3e/f: Date columns
    add_qa_slide(prs,
        "Add 'birthdate' (date) and 'last_login' (timestamp with time zone) columns.",
        "Use ALTER TABLE ADD COLUMN with appropriate data types.",
        """-- Add birthdate column
ALTER TABLE sample_table
ADD COLUMN birthdate DATE;

-- Add last_login column
ALTER TABLE sample_table
ADD COLUMN last_login TIMESTAMP WITH TIME ZONE;

-- Populate with sample data
UPDATE sample_table
SET birthdate = '2000-01-15',
    last_login = CURRENT_TIMESTAMP
WHERE id = 1;

-- View data
SELECT id, name, birthdate, last_login
FROM sample_table;"""
    )

    # Q3g: Sequence
    add_qa_slide(prs,
        "Create sequence 'sample_sequence' starting at 100, incrementing by 10.",
        "Use CREATE SEQUENCE with START WITH and INCREMENT BY clauses.",
        """-- Create sequence
CREATE SEQUENCE sample_sequence
    START WITH 100
    INCREMENT BY 10;

-- Get next value
SELECT nextval('sample_sequence');  -- Returns 100
SELECT nextval('sample_sequence');  -- Returns 110

-- Get current value (no increment)
SELECT currval('sample_sequence');  -- Returns 110

-- Set to specific value
SELECT setval('sample_sequence', 500);

-- Use in table
CREATE TABLE orders (
    id INTEGER DEFAULT nextval('sample_sequence'),
    product VARCHAR(100)
);"""
    )

    # Q3h/i: Arrays
    add_qa_slide(prs,
        "Add 'interests' array column and populate with integers.",
        "Use INTEGER[] for array type, update with ARRAY[...] or '{...}'::INTEGER[].",
        """-- Add array column
ALTER TABLE sample_table
ADD COLUMN interests INTEGER[];

-- Update with arrays
UPDATE sample_table
SET interests = ARRAY[1,3,5]
WHERE name = 'John Smith';

UPDATE sample_table
SET interests = '{2,4}'::INTEGER[]
WHERE name = 'Jane Doe';

UPDATE sample_table
SET interests = '{1,2,3,4,5}'::INTEGER[]
WHERE name = 'Bob Johnson';

-- View data
SELECT name, interests FROM sample_table;"""
    )

    # Q3j: Array contains
    add_qa_slide(prs,
        "Select rows where 'interests' contains values 1 AND 5.",
        "Use @> operator (array contains).",
        """-- Using @> operator (best method)
SELECT * FROM sample_table
WHERE interests @> ARRAY[1, 5];
-- Returns: John Smith and Bob Johnson

-- Alternative: Using ANY
SELECT * FROM sample_table
WHERE 1 = ANY(interests)
  AND 5 = ANY(interests);

-- Check overlap (&&)
SELECT * FROM sample_table
WHERE interests && ARRAY[1, 5];
-- Returns rows with either 1 OR 5"""
    )

    # Q3k: Array has but not
    add_qa_slide(prs,
        "Select rows where 'interests' contains 2 but NOT 4.",
        "Use @> for contains and NOT for negation.",
        """-- Has 2 but not 4
SELECT * FROM sample_table
WHERE interests @> ARRAY[2]
  AND NOT (interests @> ARRAY[4]);

-- Alternative method
SELECT * FROM sample_table
WHERE 2 = ANY(interests)
  AND NOT (4 = ANY(interests));

-- Using array_position
SELECT * FROM sample_table
WHERE array_position(interests, 2) IS NOT NULL
  AND array_position(interests, 4) IS NULL;"""
    )

    # Q3l: Update array
    add_qa_slide(prs,
        "Update 'interests' for id=1 to add value 2 at the end.",
        "Use array_append() or || operator.",
        """-- Using array_append
UPDATE sample_table
SET interests = array_append(interests, 2)
WHERE id = 1;

-- Using || operator
UPDATE sample_table
SET interests = interests || 2
WHERE id = 1;

-- Add multiple values
UPDATE sample_table
SET interests = interests || ARRAY[6, 7, 8]
WHERE id = 1;

-- Remove value
UPDATE sample_table
SET interests = array_remove(interests, 3)
WHERE id = 1;"""
    )

    # ========================================================================
    # SECTION 4: BACKUP & RESTORE
    # ========================================================================
    add_title_slide(prs, "Section 4", "Backup & Restore")

    # Q4a: Backup to SQL
    add_code_slide(prs,
        "Q4a: Backup database to plain SQL file",
        """# Basic backup
pg_dump mydb > mydb_backup.sql

# With username
pg_dump -U postgres mydb > mydb_backup.sql

# With host and port
pg_dump -h localhost -p 5432 -U postgres mydb > mydb_backup.sql

# Include CREATE DATABASE statement
pg_dump -U postgres -C mydb > mydb_backup.sql

# With compression
pg_dump -U postgres mydb | gzip > mydb_backup.sql.gz"""
    )

    # Q4b: Restore from SQL
    add_code_slide(prs,
        "Q4b: Restore from plain SQL file",
        """# Using psql
psql -U postgres -d mydb < mydb_backup.sql

# If database doesn't exist
createdb -U postgres mydb
psql -U postgres -d mydb < mydb_backup.sql

# From compressed backup
gunzip -c mydb_backup.sql.gz | psql -U postgres -d mydb

# Clean restore
dropdb -U postgres mydb
createdb -U postgres mydb
psql -U postgres -d mydb < mydb_backup.sql"""
    )

    # Q4c: Restore from directory
    add_code_slide(prs,
        "Q4c: Restore from directory format backup",
        """# Create directory format backup (reference)
pg_dump -U postgres -F d -f mydb_backup_dir mydb

# Restore from directory
pg_restore -U postgres -d mydb mydb_backup_dir

# Clean restore
dropdb -U postgres mydb
createdb -U postgres mydb
pg_restore -U postgres -d mydb mydb_backup_dir

# With parallelism (faster)
pg_restore -U postgres -d mydb -j 4 mydb_backup_dir

# Restore specific table only
pg_restore -U postgres -d mydb -t sample_table mydb_backup_dir"""
    )

    # Backup formats summary
    add_content_slide(prs, "Backup Formats Summary", [
        "Plain SQL (-F p): Human-readable, restore with psql",
        "Custom (-F c): Compressed binary, restore with pg_restore",
        "Directory (-F d): Parallel dump/restore, use pg_restore",
        "Tar (-F t): Archived format, use pg_restore",
        "",
        "Key flags:",
        "-C: Include CREATE DATABASE",
        "-j N: Use N parallel jobs (directory format only)",
        "--schema-only: Backup structure without data",
        "--data-only: Backup data without structure"
    ])

    # ========================================================================
    # QUICK REFERENCE SLIDES
    # ========================================================================
    add_title_slide(prs, "Quick Reference", "Essential Commands & Operations")

    # Connection management
    add_code_slide(prs,
        "Connection Management Cheat Sheet",
        """-- View all connections
SELECT pid, usename, datname, state, query
FROM pg_stat_activity;

-- Cancel query (gentle)
SELECT pg_cancel_backend(pid);

-- Kill connection (forceful)
SELECT pg_terminate_backend(pid);

-- Kill all connections to database
SELECT pg_terminate_backend(pid)
FROM pg_stat_activity
WHERE datname = 'target_db'
  AND pid <> pg_backend_pid();

-- Get your own PID
SELECT pg_backend_pid();"""
    )

    # User & permission commands
    add_code_slide(prs,
        "User & Permission Commands",
        """-- Create user with expiration
CREATE USER username WITH PASSWORD 'pass'
VALID UNTIL '2025-12-31';

-- Create role
CREATE ROLE rolename WITH INHERIT;

-- Grant privileges
GRANT SELECT, INSERT, UPDATE, DELETE
ON tablename TO username;

-- Grant on all tables
GRANT ALL ON ALL TABLES IN SCHEMA public TO role;

-- Future tables
ALTER DEFAULT PRIVILEGES IN SCHEMA public
GRANT ALL ON TABLES TO role;

-- Add user to role
GRANT rolename TO username;

-- Test as user
SET ROLE username;
RESET ROLE;"""
    )

    # Data types reference
    add_content_slide(prs, "PostgreSQL Data Types", [
        "SERIAL - Auto-incrementing integer (1, 2, 3...)",
        "INTEGER - Whole numbers (-2147483648 to 2147483647)",
        "BIGINT - Large integers",
        "NUMERIC(p,s) - Exact decimal (p=precision, s=scale)",
        "VARCHAR(n) - Variable text, max n characters",
        "TEXT - Unlimited text",
        "DATE - Date only (no time)",
        "TIMESTAMP WITH TIME ZONE - Date + time + timezone",
        "BOOLEAN - true/false",
        "INTEGER[] - Array of integers",
        "TEXT[] - Array of text"
    ])

    # Array operations
    add_code_slide(prs,
        "Array Operations Cheat Sheet",
        """-- Create array
interests INTEGER[]

-- Insert array
INSERT INTO table (arr) VALUES (ARRAY[1,2,3]);
INSERT INTO table (arr) VALUES ('{1,2,3}'::INTEGER[]);

-- Append element
UPDATE table SET arr = array_append(arr, 4);
UPDATE table SET arr = arr || 4;

-- Append array
UPDATE table SET arr = arr || ARRAY[5,6,7];

-- Check contains (has ALL of these)
WHERE arr @> ARRAY[1,5]

-- Check overlap (has ANY of these)
WHERE arr && ARRAY[1,5]

-- Access elements (1-indexed!)
arr[1]  -- First element
arr[2:4]  -- Slice

-- Array functions
array_length(arr, 1)  -- Count elements
array_position(arr, 5)  -- Find position of 5
array_remove(arr, 3)  -- Remove all 3's"""
    )

    # String functions
    add_code_slide(prs,
        "String Functions Cheat Sheet",
        """-- Padding
LPAD(string, length, fill)  -- Pad left
RPAD(string, length, fill)  -- Pad right
LPAD('hi', 5, '*')  -- '***hi'

-- Trimming
TRIM(string)   -- Remove leading & trailing spaces
LTRIM(string)  -- Remove leading spaces
RTRIM(string)  -- Remove trailing spaces

-- Case conversion
UPPER(string)  -- TO UPPERCASE
LOWER(string)  -- to lowercase

-- Substring
SUBSTRING(string, start, length)
SUBSTRING('hello', 2, 3)  -- 'ell'

-- Length
LENGTH(string)  -- Character count

-- Concatenation
string1 || string2  -- Combine strings
CONCAT(str1, str2, str3)"""
    )

    # Common queries
    add_code_slide(prs,
        "Common Query Patterns",
        """-- View table structure
\\d tablename

-- List all tables
\\dt

-- List all databases
\\l

-- Connect to database
\\c database_name

-- View current user
SELECT current_user;

-- View session user
SELECT session_user;

-- Get current date/time
SELECT CURRENT_DATE;
SELECT CURRENT_TIMESTAMP;
SELECT NOW();

-- Calculate age from birthdate
SELECT AGE(birthdate) FROM table;

-- Extract year/month/day
SELECT EXTRACT(YEAR FROM date_column);"""
    )

    # Exam tips
    add_content_slide(prs, "Exam Success Tips", [
        "✓ Remember: Arrays are 1-indexed (not 0!)",
        "✓ Always use TIMESTAMP WITH TIME ZONE",
        "✓ Use @> for 'contains all' array queries",
        "✓ pg_stat_activity shows connections",
        "✓ Never kill your own connection: pid <> pg_backend_pid()",
        "✓ SERIAL auto-increments, don't insert values",
        "✓ \\copy (client-side) vs COPY (server-side)",
        "✓ Plain SQL → psql, Binary formats → pg_restore",
        "✓ INHERIT role = members get permissions automatically",
        "✓ Test your queries before writing final answers!"
    ])

    # Save the presentation
    filename = "adv_db_midterm_study_guide.pptx"
    prs.save(filename)
    print(f"\n✓ Successfully created: {filename}")
    print(f"✓ Total slides: {len(prs.slides)}")
    return filename

if __name__ == "__main__":
    create_comprehensive_powerpoint()
