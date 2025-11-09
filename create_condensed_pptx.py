#!/usr/bin/env python3
"""
Create ultra-condensed 5-slide PowerPoint with ALL midterm answers and assignment code.
Small text for maximum information density.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_dense_slide(prs, title, content_text, font_size=7):
    """Add a slide with very dense text content."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.4))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(16)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 153)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.7), Inches(9.4), Inches(6.8))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.margin_top = 0
    content_frame.margin_bottom = 0

    content_para = content_frame.paragraphs[0]
    content_para.text = content_text
    content_para.font.name = 'Courier New'
    content_para.font.size = Pt(font_size)
    content_para.line_spacing = 0.9

    return slide

def create_ultra_condensed_pptx():
    """Create 5-slide PowerPoint with everything."""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========================================================================
    # SLIDE 1: Connections & Users/Roles
    # ========================================================================
    slide1_content = """=== SECTION 1: DATABASE CONNECTIONS ===
Q1a: List connections and PIDs?
A: SELECT pid,usename,datname,state,query_start,query FROM pg_stat_activity;
   SELECT pid,usename,datname FROM pg_stat_activity WHERE state='active';
   SELECT pid,usename FROM pg_stat_activity WHERE datname='mydb';

Q1b: Cancel/terminate connections?
A: SELECT pg_cancel_backend(12345);  -- Cancel query (gentle, SIGINT)
   SELECT pg_terminate_backend(12345);  -- Kill connection (force, SIGTERM)
   SELECT pid FROM pg_stat_activity WHERE pid=12345;  -- Verify gone

Q1c: Kill all connections for role?
A: SELECT pg_terminate_backend(pid) FROM pg_stat_activity
   WHERE usename='dev_user' AND pid<>pg_backend_pid();
   SELECT pg_terminate_backend(pid) FROM pg_stat_activity
   WHERE datname='sample_db' AND pid<>pg_backend_pid();

=== SECTION 2: USERS, ROLES & PERMISSIONS ===
Q2a: Create user with expiration, grant read/write on employee_data?
A: CREATE USER dev_user WITH PASSWORD 'SecurePass123!' VALID UNTIL '2025-12-31';
   GRANT SELECT ON employee_data TO dev_user;
   GRANT INSERT,UPDATE,DELETE ON employee_data TO dev_user;
   -- OR: GRANT SELECT,INSERT,UPDATE,DELETE ON employee_data TO dev_user;
   SELECT usename,valuntil FROM pg_user WHERE usename='dev_user';

Q2b: Create role dev_team, assign dev_user, grant all tables with inheritance?
A: CREATE ROLE dev_team WITH INHERIT;
   GRANT SELECT,INSERT,UPDATE,DELETE ON ALL TABLES IN SCHEMA public TO dev_team;
   ALTER DEFAULT PRIVILEGES IN SCHEMA public GRANT SELECT,INSERT,UPDATE,DELETE ON TABLES TO dev_team;
   GRANT dev_team TO dev_user;
   SELECT r.rolname,m.rolname FROM pg_roles r
   JOIN pg_auth_members ON r.oid=pg_auth_members.roleid
   JOIN pg_roles m ON m.oid=pg_auth_members.member WHERE r.rolname='dev_team';

Q2c: Add new developer to dev_team and test access?
A: CREATE USER new_developer WITH PASSWORD 'DevPass456!' LOGIN;
   GRANT dev_team TO new_developer;
   SET ROLE new_developer;
   SELECT * FROM employee_data LIMIT 5;
   INSERT INTO employee_data (name,email) VALUES ('Test','test@ex.com');
   UPDATE employee_data SET department='DevOps' WHERE name='Test';
   DELETE FROM employee_data WHERE name='Test';
   RESET ROLE;"""

    add_dense_slide(prs, "Slide 1: Connections & Users/Roles", slide1_content, 7)

    # ========================================================================
    # SLIDE 2: Tables, Data Types, Arrays
    # ========================================================================
    slide2_content = """=== SECTION 3: TABLES, DATA TYPES & ARRAYS ===
Q3a: Create sample_db with sample_table(id serial,name varchar(50),age numeric(4,2),description text)?
A: CREATE DATABASE sample_db;
   \\c sample_db
   CREATE TABLE sample_table(id SERIAL PRIMARY KEY,name VARCHAR(50),age NUMERIC(4,2),description TEXT);
   \\d sample_table

Q3b: Load from CSV (John Smith 25, Jane Doe 33, Bob Johnson 45)?
A: \\copy sample_table(name,age,description) FROM 'sample_data.csv' WITH CSV HEADER
   -- OR: COPY sample_table(name,age,description) FROM '/path/to/file.csv' WITH (FORMAT csv,HEADER true);
   INSERT INTO sample_table(name,age,description) VALUES
   ('John Smith',25,'Lorem ipsum'),('Jane Doe',33,'Nullam imperdiet'),('Bob Johnson',45,'Pellentesque');

Q3c: Pad name to 20 chars (left)?
A: SELECT id,LPAD(name,20,' ') AS padded_name,age,description FROM sample_table;
   SELECT RPAD(name,20,' ') FROM sample_table;  -- Right pad

Q3d: Trim whitespace from description?
A: SELECT id,name,age,TRIM(description) AS trimmed,LTRIM(description),RTRIM(description) FROM sample_table;

Q3e: Add birthdate DATE column?
A: ALTER TABLE sample_table ADD COLUMN birthdate DATE;

Q3f: Add last_login TIMESTAMP WITH TIME ZONE column?
A: ALTER TABLE sample_table ADD COLUMN last_login TIMESTAMP WITH TIME ZONE;
   UPDATE sample_table SET birthdate='2000-01-15',last_login=CURRENT_TIMESTAMP WHERE id=1;

Q3g: Create sequence starting 100, increment 10?
A: CREATE SEQUENCE sample_sequence START WITH 100 INCREMENT BY 10;
   SELECT nextval('sample_sequence');  -- Returns 100,110,120...
   SELECT currval('sample_sequence');  -- Get current (no increment)
   SELECT setval('sample_sequence',500);  -- Set to 500

Q3h: Add interests INTEGER[] column?
A: ALTER TABLE sample_table ADD COLUMN interests INTEGER[];

Q3i: Update interests: John{1,3,5}, Jane{2,4}, Bob{1,2,3,4,5}?
A: UPDATE sample_table SET interests=ARRAY[1,3,5] WHERE name='John Smith';
   UPDATE sample_table SET interests='{2,4}'::INTEGER[] WHERE name='Jane Doe';
   UPDATE sample_table SET interests='{1,2,3,4,5}'::INTEGER[] WHERE name='Bob Johnson';

Q3j: Select rows with interests containing 1 AND 5?
A: SELECT * FROM sample_table WHERE interests @> ARRAY[1,5];  -- Best: @> = contains
   SELECT * FROM sample_table WHERE 1=ANY(interests) AND 5=ANY(interests);

Q3k: Select rows with 2 but NOT 4?
A: SELECT * FROM sample_table WHERE interests @> ARRAY[2] AND NOT(interests @> ARRAY[4]);
   SELECT * FROM sample_table WHERE 2=ANY(interests) AND NOT(4=ANY(interests));

Q3l: Add 2 to end of interests for id=1?
A: UPDATE sample_table SET interests=array_append(interests,2) WHERE id=1;
   UPDATE sample_table SET interests=interests||2 WHERE id=1;
   UPDATE sample_table SET interests=interests||ARRAY[6,7,8] WHERE id=1;  -- Add multiple
   UPDATE sample_table SET interests=array_remove(interests,3) WHERE id=1;  -- Remove 3"""

    add_dense_slide(prs, "Slide 2: Tables, Data Types & Arrays", slide2_content, 6.5)

    # ========================================================================
    # SLIDE 3: Backup/Restore & Quick Reference
    # ========================================================================
    slide3_content = """=== SECTION 4: BACKUP & RESTORE ===
Q4a: Backup mydb to mydb_backup.sql?
A: pg_dump mydb > mydb_backup.sql
   pg_dump -U postgres mydb > mydb_backup.sql
   pg_dump -h localhost -p 5432 -U postgres mydb > mydb_backup.sql
   pg_dump -U postgres -C mydb > mydb_backup.sql  -- Include CREATE DATABASE
   pg_dump -U postgres mydb | gzip > mydb_backup.sql.gz  -- Compressed

Q4b: Restore from mydb_backup.sql?
A: psql -U postgres -d mydb < mydb_backup.sql
   createdb -U postgres mydb && psql -U postgres -d mydb < mydb_backup.sql
   gunzip -c mydb_backup.sql.gz | psql -U postgres -d mydb
   dropdb -U postgres mydb && createdb -U postgres mydb && psql -U postgres -d mydb < mydb_backup.sql

Q4c: Restore from mydb_backup_dir?
A: pg_dump -U postgres -F d -f mydb_backup_dir mydb  -- Create dir backup first
   pg_restore -U postgres -d mydb mydb_backup_dir
   dropdb mydb && createdb mydb && pg_restore -U postgres -d mydb mydb_backup_dir
   pg_restore -U postgres -d mydb -j 4 mydb_backup_dir  -- Parallel (4 jobs)
   pg_restore -U postgres -d mydb -t sample_table mydb_backup_dir  -- Specific table

BACKUP FORMATS: Plain SQL(-F p,default)→psql | Custom(-F c)→pg_restore | Dir(-F d)→pg_restore | Tar(-F t)→pg_restore

=== QUICK REFERENCE ===
DATA TYPES: SERIAL(auto-inc) INTEGER BIGINT NUMERIC(p,s) VARCHAR(n) TEXT DATE TIME
TIMESTAMP(use WITH TIME ZONE!) BOOLEAN INTEGER[] TEXT[] JSONB
ARRAY OPS: @>(contains) <@(contained by) &&(overlap) ||(concat) array_append(arr,val)
array_remove(arr,val) array_length(arr,1) array_position(arr,val) arr[1](1-indexed!)
STRING: LPAD(str,len,fill) RPAD TRIM LTRIM RTRIM UPPER LOWER SUBSTRING(str,start,len) LENGTH CONCAT
COMMON: \\d table \\dt \\l \\c db SELECT current_user; SELECT CURRENT_DATE; SELECT NOW();
SELECT AGE(birthdate); SELECT EXTRACT(YEAR FROM date); WHERE col BETWEEN x AND y;
GRANT SELECT,INSERT,UPDATE,DELETE ON table TO user; REVOKE SELECT ON table FROM user;
ALTER TABLE ADD COLUMN col TYPE; DROP COLUMN col; ALTER COLUMN col SET DEFAULT val;
CREATE INDEX idx ON table(col); CREATE INDEX idx ON table USING GIN(col);"""

    add_dense_slide(prs, "Slide 3: Backup/Restore & Quick Reference", slide3_content, 7)

    # ========================================================================
    # SLIDE 4: Assignment 2 Solutions
    # ========================================================================
    slide4_content = """=== ASSIGNMENT 2: DATA TYPES, ARRAYS, JSON ===
--PART 1: SERIAL,VARCHAR,TEXT,TEMPORAL TYPES--
CREATE TABLE library_books(book_id SERIAL PRIMARY KEY,title VARCHAR(100),author TEXT,
published_date DATE,available_time TIME);

INSERT INTO library_books(title,author,published_date,available_time) VALUES
('The Setup','Dan Bilzerian','2020-11-14','09:00:00'),
('1984','George Orwell','1949-06-08','10:30:00'),
('To Kill a Mockingbird','Harper Lee','1960-07-11','11:00:00'),
('The Great Gatsby','F. Scott Fitzgerald','1925-04-10','14:45:00'),
('Pride and Prejudice','Jane Austen','1813-01-28','08:15:00');

SELECT * FROM library_books WHERE published_date>'1950-1-1';
UPDATE library_books SET available_time='12:00:00' WHERE book_id=3;
SELECT * FROM library_books;

--PART 2: ARRAYS--
CREATE TABLE students(student_id SERIAL PRIMARY KEY,full_name TEXT,skills TEXT[]);

INSERT INTO students(full_name,skills) VALUES
('Hussein El Saadi',ARRAY['Data Science','Database','Programming']),
('Haidar Al Harake',ARRAY['Computing','Web Development']),
('Ahmad Taher',ARRAY['Communication','Mobile Development']);

SELECT * FROM students WHERE skills @> ARRAY['SQL'];
UPDATE students SET skills=array_append(skills,'Linear Algebra') WHERE student_id=1;
UPDATE students SET skills=array_remove(skills,'Communication') WHERE full_name='Ahmad Taher';
SELECT * FROM students;

--PART 3: JSON & JSONB--
CREATE TABLE products(product_id SERIAL PRIMARY KEY,details JSONB);

INSERT INTO products(details) VALUES
('{"name":"Smartphone","price":699.99,"tags":["electronics","gadgets","mobile"]}'),
('{"name":"Laptop","price":1299.50,"tags":["electronics","computers"]}'),
('{"name":"Wireless Earbuds","price":199.99,"tags":["electronics","audio","gadgets"]}'),
('{"name":"Apple Watch","price":549.99,"tags":["electronics","gadgets"]}');

SELECT * FROM products WHERE details->'tags'?'electronics';
UPDATE products SET details=details||'{"stock":50}'::jsonb WHERE product_id=2;
UPDATE products SET details=details-'tags' WHERE product_id=2;
SELECT * FROM products;

--JSON OPERATORS: -> (get JSON object) ->> (get text) ? (key exists) @> (contains)
|| (concatenate) - (remove key) #> (get path) #>> (get path as text)"""

    add_dense_slide(prs, "Slide 4: Assignment 2 Solutions", slide4_content, 6.5)

    # ========================================================================
    # SLIDE 5: Assignment 3 Full Text Search
    # ========================================================================
    slide5_content = """=== ASSIGNMENT 3: FULL TEXT SEARCH ===
CREATE TABLE library_books(book_id SERIAL PRIMARY KEY,title VARCHAR(200) NOT NULL,
description TEXT NOT NULL,fulltext TSVECTOR);

UPDATE library_books SET fulltext=
setweight(to_tsvector(COALESCE(title,'')),'A')||
setweight(to_tsvector(COALESCE(description,'')),'B');

CREATE TRIGGER library_books_fulltextsearch_trigger
BEFORE INSERT OR UPDATE OF title,description ON library_books
FOR EACH ROW EXECUTE PROCEDURE tsvector_update_trigger
(fulltext,'pg_catalog.english',title,description);

CREATE INDEX ix_library_books_fts_gin ON library_books USING GIN(fulltext);

INSERT INTO library_books(title,description) VALUES
('Journey to the Center of the Earth','A science fiction adventure novel by Jules Verne about a daring expedition to the earths core.'),
('The Hobbit','A fantasy novel by J.R.R. Tolkien about Bilbo Baggins quest to reclaim the Lonely Mountain.'),
('1984','A classic dystopian science fiction novel by George Orwell exploring themes of surveillance and totalitarianism.'),
('Moby Dick','A classic novel by Herman Melville about the obsessive quest of Captain Ahab to hunt the white whale.'),
('Dune','A science fiction novel by Frank Herbert exploring political intrigue and survival on the desert planet Arrakis.'),
('The Chronicles of Narnia','A fantasy series by C.S. Lewis set in the magical land of Narnia, featuring epic battles of good versus evil.'),
('Frankenstein','A classic science fiction novel by Mary Shelley about a scientist who creates a sentient creature.'),
('The Three Musketeers','A classic adventure novel by Alexandre Dumas following the exploits of a young man and three guardsmen.'),
('Twenty Thousand Leagues Under the Sea','A science fiction novel by Jules Verne detailing the underwater adventures of Captain Nemo and the Nautilus.'),
('Harry Potter and the Philosophers Stone','A fantasy novel by J.K. Rowling about a young wizards first year at Hogwarts School of Witchcraft and Wizardry.');

SELECT * FROM library_books WHERE fulltext @@ to_tsquery('fiction & adventure');
SELECT * FROM library_books WHERE fulltext @@ to_tsquery('science | fantasy');
SELECT * FROM library_books WHERE fulltext @@ to_tsquery('classic <3> novel');

--FTS FUNCTIONS: to_tsvector(text)→tsvector | to_tsquery(query)→tsquery | @@(match operator)
setweight(tsvector,weight) | ts_rank(tsvector,tsquery) | ts_headline(text,tsquery)
--TSQUERY OPS: & (AND) | (OR) ! (NOT) <N> (distance) 'word':* (prefix)
--WEIGHTS: A(highest) B C D(lowest) for ranking relevance"""

    add_dense_slide(prs, "Slide 5: Assignment 3 - Full Text Search", slide5_content, 6.5)

    # Save
    filename = "adv_db_midterm_5slides.pptx"
    prs.save(filename)
    print(f"\n✓ Created: {filename}")
    print(f"✓ Total slides: {len(prs.slides)}")
    print(f"✓ Contains: ALL midterm Q&A + Assignment 2 + Assignment 3")
    return filename

if __name__ == "__main__":
    create_ultra_condensed_pptx()
